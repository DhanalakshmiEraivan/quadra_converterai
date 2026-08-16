import os
import shutil
import subprocess
import tempfile
from pathlib import Path
from io import BytesIO

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import JSONResponse
import base64
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from starlette.background import BackgroundTask

MAX_FILE_BYTES = int(os.getenv("MAX_FILE_BYTES", str(100 * 1024 * 1024)))
ALLOWED_OFFICE = {".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".odt", ".ods", ".odp"}
ALLOWED_HTML = {".html", ".htm"}
ALLOWED_PDF = {".pdf"}

app = FastAPI(title="QuadraConverter Conversion API", version="2.0.0")
origins = [x.strip() for x in os.getenv("CORS_ORIGINS", "*").split(",") if x.strip()]
allow_credentials = "*" not in origins
app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=allow_credentials,
    allow_methods=["POST", "GET", "OPTIONS"],
    allow_headers=["*"],
)

def binary_path(*names: str) -> str:
    for name in names:
        found = shutil.which(name)
        if found:
            return found
    raise RuntimeError(f"Required converter is not installed: {', '.join(names)}")

def cleanup(path: Path):
    shutil.rmtree(path, ignore_errors=True)

def save_upload(upload: UploadFile, work: Path, allowed: set[str]) -> Path:
    suffix = Path(upload.filename or "").suffix.lower()
    if suffix not in allowed:
        raise HTTPException(400, f"Unsupported file type: {suffix or 'unknown'}")
    source = work / Path(upload.filename or f"input{suffix}").name
    size = 0
    with source.open("wb") as out:
        while True:
            chunk = upload.file.read(1024 * 1024)
            if not chunk:
                break
            size += len(chunk)
            if size > MAX_FILE_BYTES:
                raise HTTPException(413, f"File exceeds the {MAX_FILE_BYTES // (1024 * 1024)} MB conversion limit.")
            out.write(chunk)
    return source

def run_checked(args: list[str], timeout: int = 180):
    try:
        proc = subprocess.run(args, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=timeout)
    except FileNotFoundError as exc:
        raise HTTPException(500, str(exc))
    except subprocess.TimeoutExpired:
        raise HTTPException(504, "Conversion timed out. Try a smaller file.")
    if proc.returncode != 0:
        detail = (proc.stderr or proc.stdout or "Conversion failed").strip()
        raise HTTPException(422, detail[-3000:])
    return proc

def office_to_pdf(source: Path, outdir: Path, profile: Path):
    binary = binary_path("soffice", "libreoffice")
    outdir.mkdir(exist_ok=True)
    profile.mkdir(exist_ok=True)
    run_checked([
        binary, "--headless", "--convert-to", "pdf", "--outdir", str(outdir),
        "-env:UserInstallation=file://" + str(profile), str(source)
    ], timeout=240)
    pdf = outdir / f"{source.stem}.pdf"
    if not pdf.exists() or pdf.stat().st_size == 0:
        raise HTTPException(422, "LibreOffice produced no PDF.")
    return pdf

def html_to_pdf(source: Path, outdir: Path, profile: Path):
    # LibreOffice provides a deterministic server-side HTML renderer without
    # exposing the user's browser DOM or relying on client-side text extraction.
    return office_to_pdf(source, outdir, profile)

def qpdf_transform(source: Path, output: Path, password: str | None, mode: str):
    qpdf = binary_path("qpdf")
    if mode == "unlock":
        run_checked([qpdf, "--password=" + (password or ""), "--decrypt", str(source), str(output)])
    elif mode == "protect":
        if not password:
            raise HTTPException(400, "A password is required.")
        run_checked([
            qpdf, "--encrypt", password, password, "256",
            "--", str(source), str(output)
        ])
    else:
        raise HTTPException(400, "Unsupported qpdf operation.")
    if not output.exists() or output.stat().st_size == 0:
        raise HTTPException(422, "qpdf produced no output.")
    return output

def pdf_to_pdfa(source: Path, output: Path):
    # Ghostscript's PDF/A mode is used when available. This is intentionally
    # server-only because browser PDF libraries cannot produce a standards-
    # conforming PDF/A archive reliably.
    gs = binary_path("gs", "gswin64c", "gswin32c")
    run_checked([
        gs, "-dPDFA=2", "-dBATCH", "-dNOPAUSE", "-dSAFER",
        "-sDEVICE=pdfwrite",
        "-sColorConversionStrategy=RGB",
        "-sProcessColorModel=DeviceRGB",
        "-sOutputICCProfile=/usr/share/color/icc/ghostscript/srgb.icc",
        "-dPDFACompatibilityPolicy=1",
        "-o", str(output), str(source)
    ], timeout=240)
    if not output.exists() or output.stat().st_size == 0:
        raise HTTPException(422, "Ghostscript produced no PDF/A output.")
    return output

def _fitz_open(source: Path):
    import fitz
    try: return fitz.open(source)
    except Exception as exc: raise HTTPException(422, f"Could not read PDF: {exc}")

def pdf_to_docx(source: Path, output: Path):
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.section import WD_SECTION
    pdf=_fitz_open(source); doc=Document()
    for idx,page in enumerate(pdf):
        if idx: doc.add_section(WD_SECTION.NEW_PAGE)
        sec=doc.sections[-1]; rect=page.rect
        sec.page_width=Inches(max(1,rect.width/72)); sec.page_height=Inches(max(1,rect.height/72))
        sec.top_margin=sec.bottom_margin=Inches(.45); sec.left_margin=sec.right_margin=Inches(.55)
        blocks=page.get_text("dict").get("blocks",[])
        texts=[b for b in blocks if b.get("type")==0 and b.get("lines")]
        texts.sort(key=lambda b:(b["bbox"][1],b["bbox"][0]))
        for b in texts:
            para=doc.add_paragraph(); para.paragraph_format.space_after=Pt(2)
            for line in b.get("lines",[]):
                for span in line.get("spans",[]):
                    txt=span.get("text","")
                    if not txt: continue
                    r=para.add_run(txt); r.font.size=Pt(max(6,min(32,float(span.get("size",10)))))
                    if int(span.get("flags",0)) & 16: r.bold=True
                para.add_run(" ")
        for b in blocks:
            if b.get("type")!=1 or not b.get("image"): continue
            try:
                img=b["image"]; w=max(1,int(b.get("width",1))); h=max(1,int(b.get("height",1)))
                scale=min(6.7/(w/96),9/(h/96)); doc.add_paragraph().add_run().add_picture(BytesIO(img),width=Inches(max(.2,w/96*scale)))
            except Exception: pass
        if not texts and not any(b.get("type")==1 for b in blocks): doc.add_paragraph('[No extractable text on this page.]')
    doc.save(output); pdf.close(); return output

def pdf_to_pptx(source: Path, output: Path):
    from pptx import Presentation
    from pptx.util import Inches
    import fitz
    pdf=_fitz_open(source)
    if not pdf.page_count: raise HTTPException(422,'The PDF has no pages.')
    ratio=pdf[0].rect.width/max(1,pdf[0].rect.height); prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(13.333/ratio); blank=prs.slide_layouts[6]
    for page in pdf:
        pix=page.get_pixmap(matrix=fitz.Matrix(2,2),alpha=False); slide=prs.slides.add_slide(blank); slide.shapes.add_picture(BytesIO(pix.tobytes('png')),0,0,width=prs.slide_width,height=prs.slide_height)
    prs.save(output); pdf.close(); return output

def pdf_to_xlsx(source: Path, output: Path):
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font
    import pdfplumber
    wb=Workbook(); wb.remove(wb.active)
    with pdfplumber.open(str(source)) as pdf:
        for pi,page in enumerate(pdf.pages,1):
            ws=wb.create_sheet(f'Page {pi}'); wrote=False
            try:
                for table in (page.extract_tables() or []):
                    if table:
                        for row in table: ws.append([c or '' for c in row])
                        ws.append([]); wrote=True
            except Exception: pass
            words=page.extract_words(use_text_flow=True,keep_blank_chars=False) or []
            rows=[]
            for w in sorted(words,key=lambda x:(x['top'],x['x0'])):
                row=next((r for r in rows if abs(r[0]-w['top'])<=3),None)
                if row is None: row=[w['top'],[]]; rows.append(row)
                row[1].append(w)
            for _,items in rows:
                items.sort(key=lambda x:x['x0']); vals=[]; last=None
                for w in items:
                    if last is not None and w['x0']-last>24: vals.append('')
                    vals.append(w['text']); last=w['x1']
                if vals: ws.append(vals); wrote=True
            if not wrote: ws.append(['No extractable text/table data on this page.'])
            for row in ws.iter_rows():
                for cell in row: cell.alignment=Alignment(vertical='top',wrap_text=True)
            for col in ws.columns:
                letter=col[0].column_letter; mx=max((len(str(c.value or '')) for c in col),default=10); ws.column_dimensions[letter].width=min(45,max(10,mx+2))
            if ws.max_row:
                for c in ws[1]: c.font=Font(bold=True)
    wb.save(output); return output

def translation_request(text: str, target_lang: str) -> str:
    # Optional external translation service. The API is intentionally opt-in.
    import json
    import urllib.request
    endpoint = os.getenv("TRANSLATION_API_URL")
    if not endpoint:
        raise HTTPException(503, "TRANSLATION_API_URL is not configured on the conversion server.")
    payload = json.dumps({"q": text, "target": target_lang, "source": "auto", "format": "text"}).encode()
    req = urllib.request.Request(endpoint, data=payload, headers={"Content-Type": "application/json"})
    try:
        with urllib.request.urlopen(req, timeout=60) as resp:
            body = json.loads(resp.read().decode())
    except Exception as exc:
        raise HTTPException(502, f"Translation service failed: {exc}")
    return body.get("translatedText") or body.get("translation") or body.get("text") or ""

def extract_pdf_text(source: Path) -> str:
    # Used only for the optional translation tool; preserve page boundaries.
    try:
        import fitz
    except ImportError:
        raise HTTPException(500, "PyMuPDF is required for PDF translation.")
    doc = fitz.open(source)
    return "\n\n".join(page.get_text() for page in doc)
@app.post("/send-email")
async def send_email(
    file: UploadFile = File(...),
    to: str = Form(...),
    subject: str = Form(...),
    tool: str = Form(...),
):
    api_key = os.getenv("RESEND_API_KEY")
    email_from = os.getenv(
        "EMAIL_FROM",
        "QuadraConverter <onboarding@resend.dev>"
    )

    if not api_key:
        raise HTTPException(
            503,
            "Email service is not configured."
        )

    if not to or "@" not in to:
        raise HTTPException(
            400,
            "Please provide a valid email address."
        )

    content = await file.read()

    if not content:
        raise HTTPException(
            400,
            "The converted file is empty."
        )

    if len(content) > MAX_FILE_BYTES:
        raise HTTPException(
            413,
            "The file is too large to email."
        )

    try:
        import resend

        resend.api_key = api_key

        attachment = {
            "filename": file.filename or "converted-file",
            "content": list(content),
        }

        response = resend.Emails.send({
            "from": email_from,
            "to": [to],
            "subject": subject,
            "html": f"""
                <div style="font-family:Arial,sans-serif">
                    <h2>QuadraConverter</h2>

                    <p>
                        Your file has been converted successfully
                        using <strong>{tool}</strong>.
                    </p>

                    <p>
                        The converted file is attached to this email.
                    </p>

                    <p>
                        — QuadraConverter
                    </p>
                </div>
            """,
            "attachments": [attachment],
        })

        return JSONResponse({
            "success": True,
            "message": "Email sent successfully.",
            "id": response.get("id")
            if isinstance(response, dict)
            else None,
        })

    except Exception as exc:
        raise HTTPException(
            502,
            f"Email provider failed: {exc}"
        )
@app.get("/health")
def health():
    result = {"ok": True, "version": "2.0.0", "engines": []}
    for label, names in [
        ("LibreOffice", ("soffice", "libreoffice")),
        ("qpdf", ("qpdf",)),
        ("Ghostscript", ("gs", "gswin64c", "gswin32c")),
    ]:
        try:
            result["engines"].append({"name": label, "available": True, "binary": binary_path(*names)})
        except Exception:
            result["engines"].append({"name": label, "available": False})
    return result

@app.post("/convert")
async def convert(
    file: UploadFile = File(...),
    operation: str = Form(...),
    password: str = Form(""),
    targetLang: str = Form(""),
):
    work = Path(tempfile.mkdtemp(prefix="quadra-convert-"))
    try:
        suffix = Path(file.filename or "").suffix.lower()
        if operation == "office-to-pdf":
            source = save_upload(file, work, ALLOWED_OFFICE)
            pdf = office_to_pdf(source, work / "out", work / "profile")
            return FileResponse(
                pdf, media_type="application/pdf",
                filename=pdf.name,
                background=BackgroundTask(cleanup, work)
            )

        if operation == "html-to-pdf":
            source = save_upload(file, work, ALLOWED_HTML)
            pdf = html_to_pdf(source, work / "out", work / "profile")
            return FileResponse(
                pdf, media_type="application/pdf",
                filename=pdf.name,
                background=BackgroundTask(cleanup, work)
            )

        if operation in {"pdf-unlock", "pdf-protect", "pdf-to-pdfa"}:
            source = save_upload(file, work, ALLOWED_PDF)
            outdir = work / "out"
            outdir.mkdir()
            if operation == "pdf-unlock":
                output = qpdf_transform(source, outdir / f"{source.stem}-unlocked.pdf", password, "unlock")
            elif operation == "pdf-protect":
                output = qpdf_transform(source, outdir / f"{source.stem}-protected.pdf", password, "protect")
            else:
                output = pdf_to_pdfa(source, outdir / f"{source.stem}-pdfa.pdf")
            return FileResponse(
                output, media_type="application/pdf",
                filename=output.name,
                background=BackgroundTask(cleanup, work)
            )

        if operation == "pdf-to-word":
            source=save_upload(file,work,ALLOWED_PDF); outdir=work/"out"; outdir.mkdir(); output=pdf_to_docx(source,outdir/f"{source.stem}.docx")
            return FileResponse(output,media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",filename=output.name,background=BackgroundTask(cleanup,work))
        if operation == "pdf-to-pptx":
            source=save_upload(file,work,ALLOWED_PDF); outdir=work/"out"; outdir.mkdir(); output=pdf_to_pptx(source,outdir/f"{source.stem}.pptx")
            return FileResponse(output,media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",filename=output.name,background=BackgroundTask(cleanup,work))
        if operation == "pdf-to-xlsx":
            source=save_upload(file,work,ALLOWED_PDF); outdir=work/"out"; outdir.mkdir(); output=pdf_to_xlsx(source,outdir/f"{source.stem}.xlsx")
            return FileResponse(output,media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",filename=output.name,background=BackgroundTask(cleanup,work))

        if operation == "pdf-translate":
            source = save_upload(file, work, ALLOWED_PDF)
            text = extract_pdf_text(source)
            translated = translation_request(text, targetLang or "en")
            output = work / f"{source.stem}-translated.txt"
            output.write_text(translated, encoding="utf-8")
            return FileResponse(
                output, media_type="text/plain; charset=utf-8",
                filename=output.name,
                background=BackgroundTask(cleanup, work)
            )

        raise HTTPException(400, f"Unsupported conversion operation: {operation}")
    except HTTPException:
        cleanup(work)
        raise
    except Exception as exc:
        cleanup(work)
        raise HTTPException(500, str(exc))
